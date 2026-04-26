#!/usr/bin/env python3
"""
Extract a structured text inventory from a PowerPoint file.

Usage:
    python inventory.py input.pptx output.json

Output JSON structure:
    {
      "slide-0": {
        "shape-0": {
          "placeholder_type": "TITLE",
          "left": 1.5, "top": 0.5, "width": 7.0, "height": 1.2,
          "paragraphs": [
            {"text": "...", "bold": true, "font_size": 28.0, ...}
          ]
        }
      }
    }

Shapes are sorted top-to-bottom, left-to-right. Slide numbers are excluded.
"""

import argparse
import json
import sys
from copy import copy
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

EMU = 914400  # EMUs per inch
SLIDE_NUMBER_TYPE = "SLIDE_NUMBER"
ROW_TOLERANCE_IN = 0.5  # inches — shapes within this vertical distance are "same row"


# ── Paragraph extraction ──────────────────────────────────────────────────────

def _para_to_dict(para):
    d = {"text": para.text.strip()}
    if not d["text"]:
        return None

    pPr = getattr(para._p, "pPr", None)
    if pPr is not None:
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        if pPr.find(f"{ns}buChar") is not None or pPr.find(f"{ns}buAutoNum") is not None:
            d["bullet"] = True
            d["level"] = getattr(para, "level", 0)

    if para.alignment not in (None, PP_ALIGN.LEFT):
        names = {PP_ALIGN.CENTER: "CENTER", PP_ALIGN.RIGHT: "RIGHT", PP_ALIGN.JUSTIFY: "JUSTIFY"}
        if para.alignment in names:
            d["alignment"] = names[para.alignment]

    if para.space_before:
        d["space_before"] = round(para.space_before.pt, 2)
    if para.space_after:
        d["space_after"] = round(para.space_after.pt, 2)
    if para.line_spacing and hasattr(para.line_spacing, "pt"):
        d["line_spacing"] = round(para.line_spacing.pt, 2)

    if para.runs:
        run = para.runs[0]
        font = run.font
        if font.name:
            d["font_name"] = font.name
        if font.size:
            d["font_size"] = round(font.size.pt, 1)
        if font.bold is not None:
            d["bold"] = font.bold
        if font.italic is not None:
            d["italic"] = font.italic
        if font.underline is not None:
            d["underline"] = font.underline
        try:
            if font.color.rgb:
                d["color"] = str(font.color.rgb)
        except Exception:
            try:
                if font.color.theme_color:
                    d["theme_color"] = font.color.theme_color.name
            except Exception:
                pass

    return d


# ── Shape extraction ──────────────────────────────────────────────────────────

def _is_slide_number(shape):
    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
        pf = shape.placeholder_format
        if pf and pf.type and str(pf.type).split(".")[-1].split(" ")[0] == SLIDE_NUMBER_TYPE:
            return True
    return False


def _collect_shapes(shape, parent_left=0, parent_top=0):
    """Recursively collect (shape, abs_left, abs_top) for shapes with text."""
    if hasattr(shape, "shapes"):  # GroupShape
        gl = (shape.left or 0) + parent_left
        gt = (shape.top or 0) + parent_top
        result = []
        for child in shape.shapes:
            result.extend(_collect_shapes(child, gl, gt))
        return result

    if _is_slide_number(shape):
        return []
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return []
    if not shape.text_frame.text.strip():
        return []

    abs_left = (shape.left or 0) + parent_left
    abs_top = (shape.top or 0) + parent_top
    return [(shape, abs_left, abs_top)]


def _shape_to_dict(shape, abs_left, abs_top, slide):
    left_in = round(abs_left / EMU, 2)
    top_in = round(abs_top / EMU, 2)
    width_in = round((shape.width or 0) / EMU, 2)
    height_in = round((shape.height or 0) / EMU, 2)

    d = {"left": left_in, "top": top_in, "width": width_in, "height": height_in}

    if hasattr(shape, "is_placeholder") and shape.is_placeholder and shape.placeholder_format:
        pt = str(shape.placeholder_format.type).split(".")[-1].split(" ")[0]
        if pt != SLIDE_NUMBER_TYPE:
            d["placeholder_type"] = pt
            # Try to get default font size from layout
            try:
                for lp in slide.slide_layout.placeholders:
                    if lp.placeholder_format.type == shape.placeholder_format.type:
                        for elem in lp.element.iter():
                            if "defRPr" in elem.tag and elem.get("sz"):
                                d["default_font_size"] = float(elem.get("sz")) / 100
                        break
            except Exception:
                pass

    paras = [_para_to_dict(p) for p in shape.text_frame.paragraphs]
    d["paragraphs"] = [p for p in paras if p]
    return d


def _sort_shapes(shapes):
    """Sort (shape, abs_left, abs_top) by visual position: top-to-bottom, left-to-right."""
    if not shapes:
        return shapes
    by_top = sorted(shapes, key=lambda x: (x[2], x[1]))
    result, row, row_top = [], [by_top[0]], by_top[0][2] / EMU
    for item in by_top[1:]:
        item_top = item[2] / EMU
        if abs(item_top - row_top) <= ROW_TOLERANCE_IN:
            row.append(item)
        else:
            result.extend(sorted(row, key=lambda x: x[1]))
            row, row_top = [item], item_top
    result.extend(sorted(row, key=lambda x: x[1]))
    return result


# ── Main extraction ───────────────────────────────────────────────────────────

def extract_inventory(pptx_path):
    prs = Presentation(str(pptx_path))
    inventory = {}

    for slide_idx, slide in enumerate(prs.slides):
        raw = []
        for shape in slide.shapes:
            raw.extend(_collect_shapes(shape))

        if not raw:
            continue

        sorted_shapes = _sort_shapes(raw)
        slide_dict = {}
        for i, (shape, al, at) in enumerate(sorted_shapes):
            slide_dict[f"shape-{i}"] = _shape_to_dict(shape, al, at, slide)

        if slide_dict:
            inventory[f"slide-{slide_idx}"] = slide_dict

    return inventory


def main():
    parser = argparse.ArgumentParser(description="Extract text inventory from a .pptx file.")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("output", help="Output .json file")
    args = parser.parse_args()

    path = Path(args.input)
    if not path.exists():
        print(f"Error: {args.input} not found")
        sys.exit(1)

    inventory = extract_inventory(path)
    out = Path(args.output)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(inventory, indent=2, ensure_ascii=False), encoding="utf-8")

    total_shapes = sum(len(v) for v in inventory.values())
    print(f"Extracted {total_shapes} text shapes across {len(inventory)} slides → {args.output}")


if __name__ == "__main__":
    main()
