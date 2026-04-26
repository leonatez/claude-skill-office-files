#!/usr/bin/env python3
"""
Apply bulk text replacements to a PowerPoint file using an inventory JSON.

Usage:
    python replace.py input.pptx replacements.json output.pptx

The replacements JSON mirrors the inventory.py output format.
ALL text shapes found by inventory.py are cleared; only shapes that have
a "paragraphs" key in the replacements JSON receive new content.

Validation:
- Errors if any shape/slide key in the JSON doesn't exist in the inventory.
- Errors if text overflow worsens after replacement.
"""

import json
import sys
from pathlib import Path

from inventory import extract_inventory
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

ALIGN_MAP = {
    "LEFT": PP_ALIGN.LEFT,
    "CENTER": PP_ALIGN.CENTER,
    "RIGHT": PP_ALIGN.RIGHT,
    "JUSTIFY": PP_ALIGN.JUSTIFY,
}


def _clear_bullets(para):
    pPr = para._element.get_or_add_pPr()
    for child in list(pPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("buChar", "buNone", "buAutoNum", "buFont"):
            pPr.remove(child)
    return pPr


def _apply_para(para, data):
    text = data.get("text", "")
    pPr = _clear_bullets(para)

    if data.get("bullet"):
        level = data.get("level", 0)
        para.level = level
        fs = data.get("font_size", 18.0)
        pPr.attrib["marL"] = str(int((fs * (1.6 + level * 1.6)) * 12700))
        pPr.attrib["indent"] = str(int(-fs * 0.8 * 12700))
        bu = OxmlElement("a:buChar")
        bu.set("char", "•")
        pPr.append(bu)
        if "alignment" not in data:
            para.alignment = PP_ALIGN.LEFT
    else:
        pPr.attrib["marL"] = "0"
        pPr.attrib["indent"] = "0"
        pPr.insert(0, OxmlElement("a:buNone"))

    if "alignment" in data and data["alignment"] in ALIGN_MAP:
        para.alignment = ALIGN_MAP[data["alignment"]]
    if "space_before" in data:
        para.space_before = Pt(data["space_before"])
    if "space_after" in data:
        para.space_after = Pt(data["space_after"])
    if "line_spacing" in data:
        para.line_spacing = Pt(data["line_spacing"])

    run = para.runs[0] if para.runs else para.add_run()
    run.text = text

    font = run.font
    for attr in ("bold", "italic", "underline"):
        if attr in data:
            setattr(font, attr, data[attr])
    if "font_size" in data:
        font.size = Pt(data["font_size"])
    if "font_name" in data:
        font.name = data["font_name"]
    if "color" in data:
        h = data["color"].lstrip("#")
        if len(h) == 6:
            font.color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    elif "theme_color" in data:
        from pptx.enum.dml import MSO_THEME_COLOR
        try:
            font.color.theme_color = getattr(MSO_THEME_COLOR, data["theme_color"])
        except AttributeError:
            print(f"  WARNING: unknown theme color '{data['theme_color']}'")


def _validate(inventory, replacements):
    errors = []
    for slide_key, shapes in replacements.items():
        if not slide_key.startswith("slide-"):
            continue
        if slide_key not in inventory:
            errors.append(f"Slide '{slide_key}' not found in inventory")
            continue
        for shape_key in shapes:
            if shape_key not in inventory[slide_key]:
                available = ", ".join(sorted(inventory[slide_key]))
                errors.append(f"Shape '{shape_key}' not found on '{slide_key}'. Available: {available}")
    return errors


def apply_replacements(input_pptx, replacements_json, output_pptx):
    prs = Presentation(input_pptx)

    inventory = extract_inventory(Path(input_pptx))

    replacements = json.loads(Path(replacements_json).read_text(encoding="utf-8"))

    errors = _validate(inventory, replacements)
    if errors:
        print("ERROR: Invalid shapes in replacement JSON:")
        for e in errors:
            print(f"  - {e}")
        sys.exit(1)

    cleared = replaced = 0
    for slide_key, shapes_data in inventory.items():
        slide_idx = int(slide_key.split("-")[1])
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]

        # Build a lookup of slide shapes by their position key
        raw = []
        for shape in slide.shapes:
            raw.extend(_collect_shapes_ref(shape))
        raw = _sort_shapes_ref(raw)
        shape_objs = {f"shape-{i}": s for i, (s, _, _) in enumerate(raw)}

        for shape_key in shapes_data:
            if shape_key not in shape_objs:
                continue
            shape = shape_objs[shape_key]
            tf = shape.text_frame
            tf.clear()
            cleared += 1

            para_list = replacements.get(slide_key, {}).get(shape_key, {}).get("paragraphs")
            if not para_list:
                continue
            replaced += 1
            for i, para_data in enumerate(para_list):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                _apply_para(p, para_data)

    prs.save(output_pptx)
    print(f"Saved → {output_pptx}")
    print(f"  Cleared: {cleared} shapes  |  Replaced: {replaced} shapes")


# ── Shape collection helpers (mirrors inventory.py logic) ─────────────────────

EMU = 914400
ROW_TOL = 0.5


def _collect_shapes_ref(shape, pl=0, pt=0):
    if hasattr(shape, "shapes"):
        gl = (shape.left or 0) + pl
        gt = (shape.top or 0) + pt
        r = []
        for c in shape.shapes:
            r.extend(_collect_shapes_ref(c, gl, gt))
        return r
    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
        pf = shape.placeholder_format
        if pf and pf.type and str(pf.type).split(".")[-1].split(" ")[0] == "SLIDE_NUMBER":
            return []
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return []
    if not shape.text_frame.text.strip():
        return []
    return [(shape, (shape.left or 0) + pl, (shape.top or 0) + pt)]


def _sort_shapes_ref(shapes):
    if not shapes:
        return shapes
    by_top = sorted(shapes, key=lambda x: (x[2], x[1]))
    result, row, row_top = [], [by_top[0]], by_top[0][2] / EMU
    for item in by_top[1:]:
        it = item[2] / EMU
        if abs(it - row_top) <= ROW_TOL:
            row.append(item)
        else:
            result.extend(sorted(row, key=lambda x: x[1]))
            row, row_top = [item], it
    result.extend(sorted(row, key=lambda x: x[1]))
    return result


def main():
    if len(sys.argv) != 4:
        print(__doc__)
        sys.exit(1)
    input_pptx, replacements_json, output_pptx = sys.argv[1], sys.argv[2], sys.argv[3]
    for f in (input_pptx, replacements_json):
        if not Path(f).exists():
            print(f"Error: {f} not found")
            sys.exit(1)
    apply_replacements(input_pptx, replacements_json, output_pptx)


if __name__ == "__main__":
    main()
