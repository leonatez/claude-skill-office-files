#!/usr/bin/env python3
"""
Rearrange PowerPoint slides by index sequence. Slides can be repeated (duplicated).

Usage:
    python rearrange.py template.pptx output.pptx 0,34,34,50,52

Indices are 0-based. The same index may appear multiple times to duplicate a slide.
"""

import argparse
import shutil
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation


def _duplicate_slide(prs, index):
    """Append a deep copy of slide[index] to the presentation."""
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Collect image/media relationships from source
    image_rels = {
        rId: rel
        for rId, rel in source.part.rels.items()
        if "image" in rel.reltype or "media" in rel.reltype
    }

    # Clear auto-generated placeholders
    for shape in list(new_slide.shapes):
        shape.element.getparent().remove(shape.element)

    # Copy all shapes
    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        # Re-map image relationship IDs
        for blip in new_el.xpath(".//a:blip[@r:embed]"):
            old_rId = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if old_rId in image_rels:
                rel = image_rels[old_rId]
                new_rId = new_slide.part.rels.get_or_add(rel.reltype, rel._target)
                blip.set(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
                    new_rId,
                )

    return new_slide


def _delete_slide(prs, index):
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def _move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    entry = sldIdLst[from_idx]
    sldIdLst.remove(entry)
    sldIdLst.insert(to_idx, entry)


def rearrange(template_path, output_path, sequence):
    shutil.copy2(template_path, output_path)
    prs = Presentation(str(output_path))
    total = len(prs.slides)

    for idx in sequence:
        if not (0 <= idx < total):
            raise ValueError(f"Slide index {idx} out of range (0–{total - 1})")

    # Step 1: duplicate slides that appear more than once
    slide_map = []
    duplicated = {}
    for i, idx in enumerate(sequence):
        count = sequence.count(idx)
        if count == 1 or idx not in duplicated:
            slide_map.append(idx)
            if count > 1:
                dups = []
                for _ in range(count - 1):
                    _duplicate_slide(prs, idx)
                    dups.append(len(prs.slides) - 1)
                duplicated[idx] = dups
                print(f"  [{i}] slide {idx} — original (+ {count-1} duplicate(s))")
            else:
                print(f"  [{i}] slide {idx}")
        else:
            dup_idx = duplicated[idx].pop(0)
            slide_map.append(dup_idx)
            print(f"  [{i}] slide {idx} — duplicate at position {dup_idx}")

    # Step 2: delete unused slides (work backwards)
    keep = set(slide_map)
    for i in range(len(prs.slides) - 1, -1, -1):
        if i not in keep:
            _delete_slide(prs, i)
            slide_map = [x - 1 if x > i else x for x in slide_map]

    # Step 3: reorder to final sequence
    for target in range(len(slide_map)):
        current = slide_map[target]
        if current != target:
            _move_slide(prs, current, target)
            for j in range(len(slide_map)):
                if slide_map[j] > current and slide_map[j] <= target:
                    slide_map[j] -= 1
                elif slide_map[j] < current and slide_map[j] >= target:
                    slide_map[j] += 1
            slide_map[target] = target

    prs.save(str(output_path))
    print(f"\nSaved → {output_path}  ({len(prs.slides)} slides)")


def main():
    parser = argparse.ArgumentParser(description=__doc__,
                                     formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("template", help="Source .pptx file")
    parser.add_argument("output", help="Output .pptx file")
    parser.add_argument("sequence", help="Comma-separated 0-based slide indices, e.g. 0,3,3,7")
    args = parser.parse_args()

    try:
        seq = [int(x.strip()) for x in args.sequence.split(",")]
    except ValueError:
        print("Error: sequence must be comma-separated integers")
        sys.exit(1)

    template = Path(args.template)
    if not template.exists():
        print(f"Error: {args.template} not found")
        sys.exit(1)

    try:
        rearrange(template, Path(args.output), seq)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
